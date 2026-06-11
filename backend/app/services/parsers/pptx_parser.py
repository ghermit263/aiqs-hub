from pptx import Presentation


def parse_pptx(file_path: str) -> list[tuple[str, str]]:
    if file_path.lower().endswith(".ppt"):
        raise ValueError("暂不支持旧版 .ppt，请先在 PowerPoint 中另存为 .pptx")
    prs = Presentation(file_path)
    segments: list[tuple[str, str]] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        texts.append(" | ".join(cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                texts.append(f"[备注] {note}")
        if texts:
            segments.append((f"幻灯片{i}", "\n".join(texts)))
    return segments
